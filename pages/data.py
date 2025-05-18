import streamlit as st # type: ignore
import pdfplumber #type: ignore
import io
import pandas as pd # type: ignore
from docx import Document # type: ignore
import uuid
from utils.chunking import Chunker
from utils.indexing import Indexer
from config import LOGO
from io import BytesIO, FileIO
import gdown
from pptx import Presentation
import textract
import tempfile

# Page settings
st.sidebar.title('Q&A Data Assistant')
# st.sidebar.write('Immediately get your data insights')

st.sidebar.header('Nhóm tác giả')
st.sidebar.markdown("""    
    Văn Hiếu Học
                    
    Phạm Minh Thy
""")

st.logo(LOGO, size='large')

# st.write('Ngôn ngữ là' + st.session_state['language'])

# User upload data
st.subheader('Tải lên dữ liệu từ máy tính')
uploaded_files = st.file_uploader(
    label='Chọn file để tải lên', 
    type=['pdf', 'docx', 'doc', 'pptx'], 
    accept_multiple_files=True
)
st.session_state.files = uploaded_files


st.subheader('Tải lên dữ liệu từ kho lưu trữ trực tuyến')
with st.form('link_upload'):
    link = st.text_input('Dán link công khai (public link) đến tệp dữ liệu của bạn')
    submit = st.form_submit_button('Hoàn tất')

if submit:
    link_file = BytesIO()
    try:
        gdown.download(url=link, output=link_file, fuzzy=True)
        st.session_state.link_file = link_file
    except:
        st.error(':worried: Không thể tải xuống dữ liệu từ đường dẫn')

# st.header('Inputthe link: ')
# input_url = st.text_input('Input your link and press Enter')


# if st.session_state.files:
#     files_name = pd.Series([x.name for x in st.session_state.files])
#     st.dataframe(
#         data=files_name,
#         hide_index=True
#     )

# all_data = ''
if st.session_state.files or st.session_state.link_file:
    all_data = ''
    for uploaded_file in st.session_state.files:
        # print(uploaded_file.type)
        # Determine file type and read accordingly

        if uploaded_file.name.endswith(".pdf"):
            with pdfplumber.open(io.BytesIO(uploaded_file.read())) as pdf:
                for page in pdf.pages:
                    all_data += f"{page.extract_text()}\n"

        elif uploaded_file.name.endswith(".docx"):
            # Extract text from DOCX
            doc = Document(io.BytesIO(uploaded_file.read()))
            docx_text = [para.text for para in doc.paragraphs if para.text]

            for text in docx_text:
                all_data += f"{text}\n"
        elif uploaded_file.name.endswith(".doc"):
            # Extract text from .doc using textract
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp_file:
                    tmp_file.write(uploaded_file.read())
                    tmp_filepath = tmp_file.name
                raw_text = textract.process(tmp_filepath).decode("utf-8")
                all_data += f"{raw_text}\n"
            except Exception as e:
                st.error(f"Lỗi khi đọc file .doc: {e}")
        elif uploaded_file.name.endswith(".pptx"):
            presentation = Presentation(io.BytesIO(uploaded_file.read()))
            slides_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text += [shape.text]
            for text in slides_text:
                all_data += f"{text}\n"

        else:
            st.error("Định dạng file không được hỗ trợ")
            
    if st.session_state.link_file:
        doc = Document(st.session_state.link_file)
        docx_text = [para.text for para in doc.paragraphs if para.text]

        for text in docx_text:
            all_data += f"{text}\n"
    

    if all_data:
        st.success('Dữ liệu đã được xử lý thành công!')
        st.session_state.data_saved_success = True

        st.subheader("Phân mảnh dữ liệu (Chunking)")
        chunker = Chunker()
        nodes = None
        try:
            nodes = chunker.chunk(text=all_data)
            st.success('Dữ liệu được chunk thành công!')
        except:
            st.error('Đã xảy ra lỗi khi chunk dữ liệu')
    
        if nodes:
            st.subheader('Định mục dữ liệu (Indexing)')
            indexer = Indexer()
            try:
                index = indexer.index(nodes)
                st.success('Định mục dữ liệu thành công!')
                save = st.button('Lưu dữ liệu')
                if save:
                    st.success('Dữ liệu đã lưu thành công!')
                    st.session_state.data_saved_success = True
                    st.session_state['index'] = index
            except:
                st.error('Đã xảy ra lỗi khi định mục')
